import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
import matplotlib.pyplot as plt
import os

def render_latex_to_png(formula, filename):
    fig, ax = plt.subplots(figsize=(4, 1))
    fig.patch.set_visible(False)
    ax.axis('off')
    ax.text(0.5, 0.5, f"${formula}$", size=20, ha="center", va="center")
    plt.savefig(filename, bbox_inches='tight', pad_inches=0.1, dpi=300, transparent=True)
    plt.close()

def main():
    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "CSL6010-Cybersecurity Major Project:\nEyes on Your Typing & CASOM Defense"
    subtitle.text = "Analysis of SNOOPFINGER Attack and Proposed Middleware\nBy: Group of 7 Members\nRepository: https://github.com/AzDevops143/CSIEEE"

    # Slide 1: Introduction (SNOOPFINGER Attack)
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Introduction: SNOOPFINGER Attack"
    tf = slide.placeholders[1].text_frame
    tf.text = "What is SNOOPFINGER?"
    p = tf.add_paragraph()
    p.text = "A novel cross-modality side-channel attack in AR/VR systems."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Exploits head movement data (which requires zero permissions) to infer hand-based typing on virtual keyboards."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Adversaries can deduce typed words by translating 3D head orientation into 2D gaze points."
    p.level = 1

    # Slide 2: Mathematical Background
    render_latex_to_png(r"\phi = \arctan\left(\frac{2(wx+yz)}{1-2(x^2+y^2)}\right)", "eq1.png")
    render_latex_to_png(r"\theta = \arcsin(2(wy - zx))", "eq2.png")
    render_latex_to_png(r"\psi = \arctan\left(\frac{2(wz+xy)}{1-2(y^2+z^2)}\right)", "eq3.png")
    
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Understanding the Attack: Mathematical Background"
    slide.shapes.add_picture("eq1.png", Inches(1), Inches(2), width=Inches(3))
    slide.shapes.add_picture("eq2.png", Inches(1), Inches(3.5), width=Inches(3))
    slide.shapes.add_picture("eq3.png", Inches(1), Inches(5), width=Inches(3))
    txBox = slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4), Inches(3))
    tf = txBox.text_frame
    tf.text = "Conversion from Quaternions to Euler Angles"
    p = tf.add_paragraph()
    p.text = "The attack relies on converting the headset's quaternion orientation (w, x, y, z) into Euler angles (Roll, Pitch, Yaw) to estimate gaze direction."
    p.level = 1

    # Slide 3: 3D Gaze to 2D Projection
    render_latex_to_png(r"r_i = \sqrt{x_i^2 + y_i^2 + z_i^2}", "eq4.png")
    render_latex_to_png(r"\theta_i = \arctan\left(\frac{z_i}{x_i}\right)", "eq5.png")
    render_latex_to_png(r"\phi_i = \arcsin\left(\frac{y_i}{r_i}\right)", "eq6.png")

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "3D Gaze Direction to 2D Points"
    slide.shapes.add_picture("eq4.png", Inches(1), Inches(2), width=Inches(3))
    slide.shapes.add_picture("eq5.png", Inches(1), Inches(3.5), width=Inches(3))
    slide.shapes.add_picture("eq6.png", Inches(1), Inches(5), width=Inches(3))
    txBox = slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4), Inches(3))
    tf = txBox.text_frame
    tf.text = "Equirectangular Projection"
    p = tf.add_paragraph()
    p.text = "3D gaze directions are mapped onto a 2D plane (the virtual keyboard) to estimate which keys the user is looking at while typing."
    p.level = 1

    # Slide 4: Gaps & Limitations
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Identified Gaps & Limitations"
    tf = slide.placeholders[1].text_frame
    tf.text = "Limitations of the SNOOPFINGER Attack:"
    p = tf.add_paragraph()
    p.text = "Relies heavily on high-fidelity, high-frequency (72 Hz) sensor data."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Requires the user to move their head rather than just their eyes."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Gaps in Mitigation:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "The paper proposes theoretical countermeasures (e.g., Adaptive Sensor Data Obfuscation) but provides no practical implementation or architectural blueprint."
    p.level = 2

    # Slide 5: Our Proposed Solution
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Our Solution: CASOM Middleware"
    tf = slide.placeholders[1].text_frame
    tf.text = "Context-Aware Sensor Obfuscation Middleware (CASOM)"
    p = tf.add_paragraph()
    p.text = "Architecture:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Sits between the AR/VR OS Sensor Manager and Background Apps."
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Monitors sensor data to detect when a user is typing (using movement thresholding)."
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Operation:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "If typing is NOT detected, passes raw data."
    p.level = 2
    p = tf.add_paragraph()
    p.text = "If typing IS detected, injects Laplacian Noise into the data sent to background applications."
    p.level = 2

    # Slide 6: Implementation Details
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Implementation details"
    tf = slide.placeholders[1].text_frame
    tf.text = "We built a complete Python Simulation:"
    p = tf.add_paragraph()
    p.text = "generator.py: Simulates human gaze points targeting specific keys with natural micro-variances."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "attacker.py: Simulates SNOOPFINGER by clustering temporal gaze points to infer the typed word."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "defender.py: The CASOM Middleware that intercepts data and applies dynamic Laplacian noise (scale = 1.5)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "main.py: Orchestrates the attack both WITH and WITHOUT the defense, providing visualizations."
    p.level = 1

    # Slide 7: Results
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Results & Validation"
    if os.path.exists("simulation_result.png"):
        slide.shapes.add_picture("simulation_result.png", Inches(0.5), Inches(1.5), width=Inches(9))
    else:
        txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(5), Inches(1))
        txBox.text_frame.text = "[simulation_result.png not found]"

    # Slide 8: Conclusion
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Conclusion"
    tf = slide.placeholders[1].text_frame
    tf.text = "Summary"
    p = tf.add_paragraph()
    p.text = "SNOOPFINGER exploits zero-permission sensors, posing a major risk."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Our CASOM implementation successfully defends against the attack by dynamically degrading data utility for background apps."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "As proven in our simulation, the attacker infers the exact word when unprotected, but infers useless gibberish when CASOM is active."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Project Repository: https://github.com/AzDevops143/CSIEEE"
    p.level = 1

    prs.save("Cybersecurity_Major_Project.pptx")
    print("Presentation generated successfully: Cybersecurity_Major_Project.pptx")

if __name__ == "__main__":
    main()
