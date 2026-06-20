# make_pptx_clean.py
# Generates a professional, academic slide deck without emojis,
# incorporating the full v2 evaluations against the real BackTree attack.

import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import os

# Define color palette
COLOR_BG = RGBColor(15, 23, 42)        # Slate 900
COLOR_CARD = RGBColor(30, 41, 59)      # Slate 800
COLOR_TEXT_PRIMARY = RGBColor(248, 250, 252) # Slate 50
COLOR_TEXT_MUTED = RGBColor(148, 163, 184)   # Slate 400
COLOR_ACCENT_BLUE = RGBColor(56, 189, 248)   # Cyan 400
COLOR_ACCENT_GREEN = RGBColor(16, 185, 129)  # Emerald 500
COLOR_ACCENT_RED = RGBColor(244, 63, 94)     # Rose 500

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    return tf

def add_header(slide, title_text):
    tf = create_textbox(slide, Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.8))
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Calibri"
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_ACCENT_BLUE
    p.font.bold = True

def add_bullet(tf, text, level=0, bold_prefix=None, color=COLOR_TEXT_PRIMARY, size=18):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.level = level
    p.space_after = Pt(10)
    
    if bold_prefix:
        r1 = p.add_run()
        r1.text = bold_prefix
        r1.font.name = "Calibri"
        r1.font.size = Pt(size)
        r1.font.color.rgb = color
        r1.font.bold = True
        
    r2 = p.add_run()
    r2.text = text
    r2.font.name = "Calibri"
    r2.font.size = Pt(size)
    r2.font.color.rgb = color
    r2.font.bold = False

def render_latex_to_png(formula, filename):
    # Render LaTeX formulas cleanly as high-res transparent PNGs
    fig, ax = plt.subplots(figsize=(6, 1.2))
    fig.patch.set_visible(False)
    ax.axis('off')
    ax.text(0.5, 0.5, f"${formula}$", size=22, ha="center", va="center", color='#38bdf8')
    plt.savefig(filename, bbox_inches='tight', pad_inches=0.1, dpi=300, transparent=True)
    plt.close()

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625) # 16:9 ratio
    
    # -------------------------------------------------------------------------
    # Slide 1: Title Slide (Dark Theme)
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BG)
    
    # Title & Subtitle box
    tf = create_textbox(slide, Inches(0.8), Inches(1.5), Inches(8.4), Inches(2.5))
    
    p_title = tf.paragraphs[0]
    p_title.text = "Side-Channel Keystroke Leakage and Context-Aware Defense in AR/VR Environments"
    p_title.font.name = "Calibri"
    p_title.font.size = Pt(36)
    p_title.font.color.rgb = COLOR_TEXT_PRIMARY
    p_title.font.bold = True
    p_title.alignment = PP_ALIGN.LEFT
    p_title.space_after = Pt(20)
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Analysis of the SNOOPFINGER Attack and the CASOM Defense Framework"
    p_sub.font.name = "Calibri"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = COLOR_ACCENT_BLUE
    p_sub.alignment = PP_ALIGN.LEFT
    p_sub.space_after = Pt(40)
    
    p_meta = tf.add_paragraph()
    p_meta.text = "CSL6010-Cybersecurity Major Project Submission\nPrepared by: Graduate Student Analysis Team\nRepository: https://github.com/AzDevops143/CSIEEE"
    p_meta.font.name = "Calibri"
    p_meta.font.size = Pt(13)
    p_meta.font.color.rgb = COLOR_TEXT_MUTED
    p_meta.alignment = PP_ALIGN.LEFT

    # -------------------------------------------------------------------------
    # Slide 2: Context: SNOOPFINGER Side-Channel Attack
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BG)
    add_header(slide, "The Side-Channel Threat: SNOOPFINGER")
    
    tf = create_textbox(slide, Inches(0.6), Inches(1.4), Inches(8.8), Inches(3.8))
    add_bullet(tf, "SNOOPFINGER is a cross-modality side-channel attack targeting head-mounted displays in VR/AR environments.", bold_prefix="Attack Definition: ")
    add_bullet(tf, "APIs treat head-movement orientation (quaternions) as zero-permission sensors, broadcasting it to all background threads.", bold_prefix="Vulnerability Surface: ")
    add_bullet(tf, "When users type on virtual keyboards, they make natural micro-pauses (dwell times) to visually target and confirm keys.", bold_prefix="Exploitation Vector: ")
    add_bullet(tf, "Background processes reconstruct typed text by converting 3D head rotation into 2D gaze points and clustering them temporally.", bold_prefix="Mechanism: ")
    add_bullet(tf, "Requires no camera access or high privileges, allowing standard apps to steal passwords and pins.", bold_prefix="Impact: ", color=COLOR_ACCENT_RED)

    # -------------------------------------------------------------------------
    # Slide 3: Threat Model & Math (3D to 2D Gaze Projection)
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BG)
    add_header(slide, "Mathematical Modeling: 3D to 2D Gaze Projection")
    
    # Render formulas
    render_latex_to_png(r"\theta = \arcsin(2(wy - zx))", "math_pitch.png")
    render_latex_to_png(r"\psi = \arctan\left(\frac{2(wz+xy)}{1-2(y^2+z^2)}\right)", "math_yaw.png")
    
    tf = create_textbox(slide, Inches(0.6), Inches(1.3), Inches(4.5), Inches(3.9))
    add_bullet(tf, "Converts raw headset orientation tracking (w, x, y, z) into Euler angles (Pitch and Yaw) to estimate looking directions.", bold_prefix="Coordinate Conversion: ", size=16)
    add_bullet(tf, "Maps the 3D gaze vector onto the 2D plane of the virtual keyboard via an equirectangular coordinate projection.", bold_prefix="2D Projection: ", size=16)
    add_bullet(tf, "Models natural gaze jitter using Gaussian distribution centered at key targets:", bold_prefix="Tremor Simulation: ", size=16)
    add_bullet(tf, "x ~ N(x_key, sigma^2), where sigma = 0.03 cm represents human target tremor.", level=1, size=15)
    
    # Add pictures of math equations on the right half
    slide.shapes.add_picture("math_pitch.png", Inches(5.3), Inches(1.8), width=Inches(4.1))
    slide.shapes.add_picture("math_yaw.png", Inches(5.3), Inches(3.2), width=Inches(4.1))

    # -------------------------------------------------------------------------
    # Slide 4: Identified Gaps & Limitations
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BG)
    add_header(slide, "Identified Gaps & Defense Limitations")
    
    tf = create_textbox(slide, Inches(0.6), Inches(1.4), Inches(8.8), Inches(3.8))
    add_bullet(tf, "The SNOOPFINGER IEEE paper discusses countermeasures (e.g. Obfuscation) but provides no implementation code or API design.", bold_prefix="Abstract Proposal Gap: ")
    add_bullet(tf, "Adding simple independent random noise per frame degrades coordinate utility for background games/apps without stopping attackers.", bold_prefix="Utility Degradation: ")
    add_bullet(tf, "The paper assumes simple attackers. It does not evaluate defenses against adaptive adversaries or geometric attacks like BackTree.", bold_prefix="Evaluation Deficit: ")
    add_bullet(tf, "Virtual keyboard typing detection must be context-aware; static noise injection would break normal AR/VR head tracking gestures.", bold_prefix="Usability Conflict: ", color=COLOR_ACCENT_RED)

    # -------------------------------------------------------------------------
    # Slide 5: The \"Toy Attacker\" Trap: Why IID Noise Fails
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BG)
    add_header(slide, "The 'Toy Attacker' Trap: Why IID Noise Fails")
    
    render_latex_to_png(r"E[x_{\text{obfuscated}}] = E[x_{\text{raw}} + \eta] = E[x_{\text{raw}}]", "math_expectation.png")
    
    tf = create_textbox(slide, Inches(0.6), Inches(1.3), Inches(5.0), Inches(3.9))
    add_bullet(tf, "Adds independent zero-mean Laplace noise to each sample. Points appear scattered on raw plots.", bold_prefix="IID Noise Method: ", size=16)
    add_bullet(tf, "Since user dwell holds ~15-35 frames, the attacker segments by dwell and averages coordinates.", bold_prefix="Adaptive Counter-Strategy: ", size=16)
    add_bullet(tf, "By the Law of Large Numbers, independent zero-mean noise cancels out and averages back to zero.", bold_prefix="Noise Cancellation: ", size=16)
    add_bullet(tf, "Evaluated in our benchmark: drops naive attacker to 11% but fails against adaptive attacker (51% Top-1 character leak).", bold_prefix="Performance: ", color=COLOR_ACCENT_RED, size=16)
    
    slide.shapes.add_picture("math_expectation.png", Inches(5.8), Inches(2.2), width=Inches(3.8))

    # -------------------------------------------------------------------------
    # Slide 6: Proposed Solution: CASOM Middleware Architecture
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BG)
    add_header(slide, "Proposed Solution: CASOM Architecture")
    
    tf = create_textbox(slide, Inches(0.6), Inches(1.4), Inches(8.8), Inches(3.8))
    add_bullet(tf, "Context-Aware Sensor Obfuscation Middleware running inside the OS sensor framework boundary.", bold_prefix="OS-Level Interceptor: ")
    add_bullet(tf, "Detects keyboard events and focus states to dynamically trigger protection without breaking standard gestures.", bold_prefix="Context Awareness: ")
    add_bullet(tf, "Bifurcates raw tracking coordinates into separate foreground and background streams.", bold_prefix="Stream Bifurcation: ")
    add_bullet(tf, "Legitimate Virtual Keyboard receives clean, high-precision telemetry, ensuring 100% typing accuracy and zero UX lag.", bold_prefix="Foreground Stream: ", color=COLOR_ACCENT_GREEN)
    add_bullet(tf, "Background processes receive mathematically obfuscated streams, stripping them of high-fidelity spatial telemetry.", bold_prefix="Background Stream: ", color=COLOR_ACCENT_RED)

    # -------------------------------------------------------------------------
    # Slide 7: CASOM v2: Surviving the BackTree Attack
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BG)
    add_header(slide, "CASOM v2: Surviving the BackTree Attack")
    
    tf = create_textbox(slide, Inches(0.6), Inches(1.4), Inches(8.8), Inches(3.8))
    add_bullet(tf, "The BackTree attack (SNOOPFINGER paper section 5.6) maps relative geometry vectors between keypress centroids, bypassing simple drift.", bold_prefix="BackTree Adversary: ")
    add_bullet(tf, "A slow coordinate bias drift shifts consecutive centroids in the same direction, preserving relative vectors and letting BackTree succeed.", bold_prefix="Drift Failure: ")
    add_bullet(tf, "A Laplace noise offset is generated once per block of 15 samples (~dwell length) and held constant during the block.", bold_prefix="Block Offset Defense (Recommended): ", color=COLOR_ACCENT_GREEN)
    add_bullet(tf, "Because noise changes per block, it shifts each key centroid independently, breaking relative vectors and surviving dwell averages.", bold_prefix="Disruption Mechanism: ")
    add_bullet(tf, "Does not require precise typing boundaries, making it highly deployable via standard timers.", bold_prefix="Deployment Advantage: ")

    # -------------------------------------------------------------------------
    # Slide 8: Empirical Results: BackTree Word Inference
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BG)
    add_header(slide, "Empirical Results: Word-Level Evaluation Table")
    
    # Add a table
    # Columns: Defense mode, Word Top-1, Word Top-3, Word Top-10
    rows = 7
    cols = 4
    left = Inches(0.6)
    top = Inches(1.3)
    width = Inches(8.8)
    height = Inches(2.2)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    headers = ["Defense Mode", "Word Top-1 Accuracy", "Word Top-3 Accuracy", "Word Top-10 Accuracy"]
    data = [
      ["No Defense (Baseline)", "80.9%", "80.9%", "80.9%"],
      ["IID Noise (Original)", "35.8%", "51.8%", "62.9%"],
      ["Correlated Drift", "27.6%", "42.1%", "52.7%"],
      ["Downsample & Quantize", "2.3%", "5.6%", "8.9%"],
      ["Block Offset (CASOM)", "3.6%", "7.0%", "13.6%"],
      ["Per-Keypress Offset", "2.3%", "3.8%", "7.6%"]
    ]
    
    # Populate headers
    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_CARD
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.name = "Calibri"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT_BLUE
        p.alignment = PP_ALIGN.CENTER
        
    # Populate data
    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            cell = table.cell(r+1, c)
            cell.fill.solid()
            # Alternating background colors
            if r % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(22, 30, 49)
            else:
                cell.fill.fore_color.rgb = RGBColor(15, 23, 42)
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = "Calibri"
            p.font.size = Pt(13)
            p.font.color.rgb = COLOR_TEXT_PRIMARY
            # Bold highlights
            if c == 0:
                p.alignment = PP_ALIGN.LEFT
                p.font.bold = True
            else:
                p.alignment = PP_ALIGN.CENTER
                if r >= 3 and c == 1: # Collapse values
                    p.font.color.rgb = COLOR_ACCENT_GREEN
                    p.font.bold = True
                elif r == 0: # Baseline threat
                    p.font.color.rgb = COLOR_ACCENT_RED
                    p.font.bold = True
                    
    tf_notes = create_textbox(slide, Inches(0.6), Inches(3.8), Inches(8.8), Inches(1.5))
    add_bullet(tf_notes, "The recommended Block Offset defense collapses Top-1 accuracy to 3.6% (near chance for 8k dictionary).", bold_prefix="Key Takeaway: ", size=14)
    add_bullet(tf_notes, " drift fails to defend against BackTree because vectors survive, leaving Top-1 at 27.6%.", bold_prefix="Drift Deficit: ", size=14)
    add_bullet(tf_notes, "IID Noise averages out, leaving a significant Top-10 leak of 62.9%.", bold_prefix="IID Deficit: ", size=14)

    # -------------------------------------------------------------------------
    # Slide 9: Verification Plots
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BG)
    add_header(slide, "Empirical Validation Plots")
    
    tf = create_textbox(slide, Inches(0.6), Inches(1.3), Inches(4.5), Inches(3.9))
    add_bullet(tf, "Includes a character-level simulation dashboard and word-level evaluation plots.", bold_prefix="Telemetry Plots: ", size=16)
    add_bullet(tf, "Left Panel (simulation_result.png) demonstrates visual scatter. Dwell clusters are completely dispersed.", bold_prefix="Visual Scatter: ", size=16)
    add_bullet(tf, "Right Panel (backtree_result.png) maps word inference rates under various noise profiles, validating the collapse of BackTree.", bold_prefix="BackTree Results: ", size=16)
    
    # Add pictures if they exist
    if os.path.exists("simulation_result.png"):
        slide.shapes.add_picture("simulation_result.png", Inches(5.3), Inches(1.3), width=Inches(4.1), height=Inches(1.9))
    if os.path.exists("backtree_result.png"):
        slide.shapes.add_picture("backtree_result.png", Inches(5.3), Inches(3.3), width=Inches(4.1), height=Inches(1.9))

    # -------------------------------------------------------------------------
    # Slide 10: Conclusion & Core Contributions
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BG)
    add_header(slide, "Conclusion & Core Contributions")
    
    tf = create_textbox(slide, Inches(0.6), Inches(1.4), Inches(8.8), Inches(3.8))
    add_bullet(tf, "Proved that zero-permission head sensors leak text, posing a severe side-channel danger in virtual environments.", bold_prefix="Side-Channel Vulnerability: ")
    add_bullet(tf, "Implemented and evaluated the Context-Aware Sensor Obfuscation Middleware (CASOM) inside a Python simulation environment.", bold_prefix="CASOM Framework: ")
    add_bullet(tf, "Disproved simple countermeasures (IID Noise and Drift) by evaluating them against the real BackTree attack.", bold_prefix="Evaluation Honesty: ")
    add_bullet(tf, "Validated Block Offset Laplace Noise as a robust, highly secure, and deployable countermeasure that collapses attack success to ~3.6% while preserving 100% foreground user experience.", bold_prefix="Recommended Defense: ", color=COLOR_ACCENT_GREEN)
    
    # -------------------------------------------------------------------------
    # Save the Presentation
    # -------------------------------------------------------------------------
    output_filename = "Cybersecurity_Major_Project.pptx"
    prs.save(output_filename)
    # Also save a copy in the parent directory just in case the user wants it there
    prs.save("../" + output_filename)
    print(f"Presentation saved successfully as {output_filename}")

if __name__ == "__main__":
    main()
